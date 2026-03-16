from pptx import Presentation
prs = Presentation("ppt format.pptx")
for idx, slide in enumerate(prs.slides, 1):
    print(f"--- Slide {idx} ---")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            print(f"  > {para.text.strip()}")